"""Document preprocessing for multimodal chat attachments.

Converts non-image file types (PDF, DOCX, PPTX, XLSX, ZIP) into structured
text and/or rendered images that can be fed to multimodal models as
OpenAI-style content parts.

Each extractor returns a list of dicts like:
    [{"type": "text", "text": "..."},
     {"type": "image_url", "image_url": {"url": "data:image/png;base64,..."}},
     ...]
"""

from __future__ import annotations

import base64
import io
import logging
import os
import zipfile

from .prompts import wrap_attachment_content

logger = logging.getLogger("bimo.document_processor")


# --- helpers ---

def _bytes_to_data_uri(png_bytes: bytes, mime: str = "image/png") -> str:
    b64 = base64.b64encode(png_bytes).decode("ascii")
    return f"data:{mime};base64,{b64}"


def _image_part(img_bytes: bytes, mime: str = "image/png") -> dict:
    return {"type": "image_url", "image_url": {"url": _bytes_to_data_uri(img_bytes, mime)}}


def _text_part(text: str) -> dict:
    return {"type": "text", "text": text}


def _ext(filename: str) -> str:
    return (os.path.splitext(filename)[1] or "").lower().lstrip(".")


def _guess_mime_from_ext(ext: str) -> str:
    mapping = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "zip": "application/zip",
        "txt": "text/plain",
        "py": "text/x-python",
        "js": "application/javascript",
        "json": "application/json",
        "csv": "text/csv",
        "md": "text/markdown",
        "html": "text/html",
        "xml": "application/xml",
    }
    return mapping.get(ext, "application/octet-stream")


# --- PDF ---

def _extract_pdf(file_bytes: bytes, max_pages: int = 30) -> list[dict]:
    """Render pages to JPEG at adaptive DPI and extract text."""
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is not installed") from exc

    _has_pil = False
    try:
        from PIL import Image
        _has_pil = True
    except ImportError:
        logger.warning("Pillow not installed; PDF pages will render as PNG (much larger).")

    _resample = None
    if _has_pil:
        try:
            _resample = Image.Resampling.LANCZOS  # Pillow 10.x
        except AttributeError:
            _resample = Image.LANCZOS  # Pillow < 10

    parts: list[dict] = []
    text_chunks: list[str] = []

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    total_pages = len(doc)
    pages_to_process = min(total_pages, max_pages)

    if total_pages > max_pages:
        text_chunks.append(f"[PDF has {total_pages} pages; showing first {max_pages}]")

    # Adaptive DPI: lighter for large page counts or large file sizes.
    file_mb = len(file_bytes) / (1024 * 1024)
    if file_mb > 15 or pages_to_process > 5:
        dpi = 120
    elif file_mb > 8 or pages_to_process > 3:
        dpi = 150
    else:
        dpi = 170
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)

    total_img_bytes = 0
    max_per_page_kb = 350  # per-page soft cap

    for page_num in range(pages_to_process):
        page = doc.load_page(page_num)
        # text
        txt = page.get_text()
        if txt.strip():
            text_chunks.append(f"--- Page {page_num + 1} ---\n{txt.strip()}")
        # image
        pix = page.get_pixmap(matrix=mat)

        if _has_pil:
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            # If page is still too large, downscale proportionally.
            scale = 1.0
            while True:
                bio = io.BytesIO()
                w, h = int(pix.width * scale), int(pix.height * scale)
                if scale < 1.0:
                    img_resized = img.resize((w, h), _resample)
                else:
                    img_resized = img
                img_resized.save(bio, format="JPEG", quality=60, optimize=True)
                img_bytes = bio.getvalue()
                if len(img_bytes) <= max_per_page_kb * 1024 or scale <= 0.5:
                    break
                scale -= 0.15
            mime = "image/jpeg"
        else:
            img_bytes = pix.tobytes("png")
            mime = "image/png"

        total_img_bytes += len(img_bytes)
        parts.append(_image_part(img_bytes, mime))
        logger.info(
            "pdf: page %d/%d dpi=%d scale=%.2f %s=%.1fKB total=%.1fMB",
            page_num + 1, pages_to_process, dpi,
            scale if _has_pil else 1.0,
            mime.split("/")[1].upper(), len(img_bytes) / 1024,
            total_img_bytes / (1024 * 1024),
        )

    doc.close()

    if text_chunks:
        parts.insert(0, _text_part("\n\n".join(text_chunks)))
    return parts


# --- DOCX ---

def _extract_docx(file_bytes: bytes) -> list[dict]:
    """Extract paragraphs and tables as markdown-ish text."""
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is not installed") from exc

    doc = Document(io.BytesIO(file_bytes))
    chunks: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            style = para.style.name.lower() if para.style else ""
            if "heading" in style:
                level = style.replace("heading", "").strip()
                try:
                    hl = max(1, min(6, int(level)))
                    chunks.append(f"{'#' * hl} {text}")
                except ValueError:
                    chunks.append(f"## {text}")
            else:
                chunks.append(text)

    # Tables
    for table in doc.tables:
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if rows:
            chunks.append("\n[Table]\n" + "\n".join(" | ".join(r) for r in rows))

    full_text = "\n\n".join(chunks)
    if not full_text.strip():
        full_text = "[Document appears to be empty or contains no extractable text]"

    return [_text_part(full_text)]


# --- XLSX ---

def _extract_xlsx(file_bytes: bytes) -> list[dict]:
    """Extract each sheet as a simple markdown table."""
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl is not installed") from exc

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    chunks: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        chunks.append(f"--- Sheet: {sheet_name} ---")
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(cell) if cell is not None else "" for cell in row])
        if rows:
            chunks.append(" | ".join(rows[0]))
            chunks.append(" | ".join(["---"] * len(rows[0])))
            for r in rows[1:]:
                chunks.append(" | ".join(r))
        chunks.append("")

    wb.close()
    full_text = "\n".join(chunks)
    if not full_text.strip():
        full_text = "[Spreadsheet appears to be empty]"

    return [_text_part(full_text)]


# --- PPTX ---

def _extract_pptx(file_bytes: bytes, max_slides: int = 10) -> list[dict]:
    """Extract text from slides (titles, bullets, notes)."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is not installed") from exc

    prs = Presentation(io.BytesIO(file_bytes))
    chunks: list[str] = []
    total_slides = len(prs.slides)
    slides_to_process = min(total_slides, max_slides)

    if total_slides > max_slides:
        chunks.append(f"[Presentation has {total_slides} slides; showing first {max_slides}]\n")

    for idx, slide in enumerate(list(prs.slides)[:slides_to_process], start=1):
        slide_chunks: list[str] = [f"--- Slide {idx} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_chunks.append(shape.text.strip())
            if shape.has_table:
                tbl = shape.table
                rows: list[list[str]] = []
                for row in tbl.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    slide_chunks.append("\n".join(" | ".join(r) for r in rows))
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_chunks.append(f"[Notes] {notes}")
        chunks.append("\n".join(slide_chunks))

    full_text = "\n\n".join(chunks)
    if not full_text.strip():
        full_text = "[Presentation appears to be empty]"

    return [_text_part(full_text)]


# --- ZIP ---

def _extract_zip(
    file_bytes: bytes,
    *,
    max_files: int = 20,
    max_total_mb: int = 50,
) -> list[dict]:
    """Extract one level deep with zip-bomb and path-traversal safeguards."""
    parts: list[dict] = []
    processed = 0
    total_extracted = 0
    max_total_bytes = max_total_mb * 1024 * 1024

    with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
        # Filter to safe files: skip dirs, reject traversal paths ('..', leading '/')
        members = [
            m for m in zf.infolist()
            if not m.is_dir()
            and ".." not in m.filename
            and not m.filename.startswith("/")
            and not m.filename.startswith("\\")
            and "/" not in m.filename.strip("/")
        ]

        for member in members:
            if processed >= max_files:
                parts.append(_text_part(f"[ZIP: {len(members)} files total; only first {max_files} processed]"))
                break

            ext = _ext(member.filename)
            if ext not in {"pdf", "docx", "pptx", "xlsx", "txt", "py", "js", "json", "csv", "md", "html", "xml", "png", "jpg", "jpeg", "gif", "webp"}:
                continue

            # Zip-bomb guard 1: check compression ratio (abnormally high expansion ratio)
            if member.compress_size > 0 and (member.file_size / member.compress_size) > 100:
                logger.warning("zip: rejected suspicious decompression ratio on %s", member.filename)
                parts.append(_text_part(f"[ZIP: skipped suspicious file with high compression ratio: {member.filename}]"))
                continue

            # Zip-bomb guard 2: check declared uncompressed size against quota
            if member.file_size > max_total_bytes or total_extracted + member.file_size > max_total_bytes:
                parts.append(_text_part(f"[ZIP: total extracted size exceeded {max_total_mb} MB; stopping]"))
                break

            try:
                buf = bytearray()
                over_quota = False
                with zf.open(member) as src:
                    while True:
                        chunk = src.read(64 * 1024)
                        if not chunk:
                            break
                        buf.extend(chunk)
                        if total_extracted + len(buf) > max_total_bytes:
                            over_quota = True
                            break
                if over_quota:
                    parts.append(_text_part(f"[ZIP: total extracted size exceeded {max_total_mb} MB; stopping]"))
                    break
                data = bytes(buf)
            except Exception as exc:
                logger.warning("zip: failed to read %s: %s", member.filename, exc)
                continue

            total_extracted += len(data)

            processed += 1
            mime = _guess_mime_from_ext(ext)
            fake_attachment = {
                "filename": member.filename,
                "content_type": mime,
                "size": len(data),
            }

            try:
                nested = _dispatch_by_type(fake_attachment, data)
                if nested:
                    parts.append(_text_part(f"\n--- File in ZIP: {member.filename} ---\n"))
                    parts.extend(nested)
            except Exception as exc:
                logger.warning("zip: failed to process %s: %s", member.filename, exc)
                parts.append(_text_part(f"[Could not process {member.filename}: {exc}]"))

    if not parts:
        return [_text_part("[ZIP archive is empty or contains no supported files]")]

    return parts


# --- generic text fallback ---

def _extract_text_file(file_bytes: bytes, filename: str) -> list[dict]:
    """Best-effort decode as UTF-8 text."""
    try:
        text = file_bytes.decode("utf-8", errors="replace")
        return [_text_part(f"--- {filename} ---\n{text}")]
    except Exception:
        return [_text_part(f"[Could not decode {filename} as text]")]


# --- dispatcher ---

def _dispatch_by_type(attachment: dict, file_bytes: bytes) -> list[dict]:
    content_type = (attachment.get("content_type") or "").lower()
    filename = attachment.get("filename", "")
    ext = _ext(filename)

    if content_type.startswith("image/") or ext in {"png", "jpg", "jpeg", "gif", "webp", "bmp"}:
        return []

    if "pdf" in content_type or ext == "pdf":
        return _extract_pdf(file_bytes)
    if "wordprocessingml" in content_type or ext == "docx":
        return _extract_docx(file_bytes)
    if "spreadsheetml" in content_type or ext == "xlsx":
        return _extract_xlsx(file_bytes)
    if "presentationml" in content_type or ext == "pptx":
        return _extract_pptx(file_bytes)
    if content_type == "application/zip" or ext == "zip":
        return _extract_zip(file_bytes)
    if ext in {"txt", "py", "js", "json", "csv", "md", "html", "xml"}:
        return _extract_text_file(file_bytes, filename)

    return []


# --- public API ---

def process_attachment(attachment: dict, download_fn) -> list[dict]:
    """Process an attachment into OpenAI-style content parts.

    ``download_fn`` is a callable that accepts the attachment's ``path``
    string and returns raw ``bytes``.
    """
    path = attachment.get("path")
    filename = attachment.get("filename") or path or "attachment"
    if not path:
        return []

    try:
        file_bytes = download_fn(path)
    except Exception as exc:
        logger.warning("document_processor: failed to download %s: %s", path, exc)
        return []

    if not file_bytes:
        return []

    try:
        parts = _dispatch_by_type(attachment, file_bytes)
        # Wrap all text parts in structured boundary delimiters to protect against prompt injection
        delimited_parts: list[dict] = []
        for p in parts:
            if p.get("type") == "text" and p.get("text"):
                delimited_parts.append({
                    "type": "text",
                    "text": wrap_attachment_content(filename, p["text"]),
                })
            else:
                delimited_parts.append(p)

        if delimited_parts:
            logger.info(
                "document_processor: %s -> %d parts",
                filename,
                len(delimited_parts),
            )
        return delimited_parts
    except Exception as exc:
        logger.warning(
            "document_processor: failed to parse %s: %s",
            filename,
            exc,
        )
        return []
